"""Generate final report (DOCX) and presentation (PPTX) for Myntra Data Analysis."""

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
REPORTS_DIR = ROOT / "reports"
IMAGES_DIR = ROOT / "images"
REPORTS_DIR.mkdir(exist_ok=True)


def add_heading(doc, text, level=1):
    doc.add_heading(text, level=level)


def add_para(doc, text, bold=False):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    return p


def add_image(doc, image_path: Path, width_inches: float = 5.5, caption: str | None = None):
    if not image_path.exists():
        return
    doc.add_picture(str(image_path), width=Inches(width_inches))
    if caption:
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER


def build_report():
    doc = Document()

    # Title page
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("\n\nMyntra Men's Jeans\nData Analysis Report\n")
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0xE6, 0x2E, 0x5C)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = subtitle.add_run(
        "\nEnd-to-End Data Analytics Project\n"
        "Pricing · Brand Performance · Customer Engagement\n\n"
        "July 2026"
    )
    sub_run.font.size = Pt(14)

    doc.add_page_break()

    # Executive Summary
    add_heading(doc, "Executive Summary")
    doc.add_paragraph(
        "This report presents a comprehensive analysis of Myntra's men's jeans product catalog, "
        "covering 31,527 cleaned product listings across 371 brands. The analysis examines pricing "
        "strategies, discount patterns, customer ratings, and brand engagement to generate actionable "
        "business recommendations for Myntra's merchandising, marketing, and partnership teams."
    )
    doc.add_paragraph(
        "Key findings indicate that the mid-price segment (₹1,001–2,000) drives the majority of "
        "customer engagement, brands such as HIGHLANDER and HARDSODA deliver the strongest overall "
        "business performance, and value-for-money leaders like MarvelQ and AngelFab combine "
        "competitive pricing with strong customer satisfaction."
    )

    # Introduction
    add_heading(doc, "1. Introduction")
    doc.add_paragraph(
        "Myntra is one of India's leading fashion e-commerce platforms, offering a wide range of "
        "apparel across price segments. Understanding how products are priced, discounted, and "
        "rated by customers is essential for optimizing inventory, promotional campaigns, and "
        "brand partnerships."
    )
    doc.add_paragraph(
        "This project follows a structured analytics pipeline: data loading, quality assessment, "
        "cleaning, exploratory data analysis (EDA), and advanced business insights. The objective "
        "is to transform raw product listing data into data-driven recommendations that support "
        "Myntra's business strategy."
    )

    # Dataset Description
    add_heading(doc, "2. Dataset Description")
    doc.add_paragraph(
        "The dataset was obtained through web scraping of Myntra's men's jeans product listings. "
        "It contains seven attributes per product record."
    )
    table = doc.add_table(rows=8, cols=3)
    table.style = "Light Grid Accent 1"
    headers = ["Column", "Type", "Description"]
    rows_data = [
        ["brand_name", "Categorical", "Brand of the product"],
        ["pants_description", "Categorical", "Product title / description"],
        ["price", "Numerical", "Selling price in INR"],
        ["MRP", "Numerical", "Maximum retail price in INR"],
        ["discount_percent", "Numerical", "Discount as decimal (0–1)"],
        ["ratings", "Numerical", "Average customer rating (1–5)"],
        ["number_of_ratings", "Numerical", "Total number of ratings"],
    ]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    for r, row in enumerate(rows_data, start=1):
        for c, val in enumerate(row):
            table.rows[r].cells[c].text = val

    doc.add_paragraph("")
    doc.add_paragraph(
        "Original dataset: 52,120 records. After cleaning: 31,527 records across 371 unique brands."
    )

    # Data Cleaning Summary
    add_heading(doc, "3. Data Cleaning Summary")
    doc.add_paragraph("The following data quality issues were identified and addressed:")
    for item in [
        "Missing values: None found — no imputation required.",
        "Duplicate rows: 17,047 exact duplicates removed (52,120 → 35,073).",
        "Inconsistent discounts: 3,546 records with discount values outside the "
        "0–1 decimal range (scraper format inconsistency) removed (35,073 → 31,527).",
        "Data types: All numerical columns validated; discount stored as decimal (0–1).",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    # EDA Highlights
    add_heading(doc, "4. EDA Highlights")
    add_image(
        doc,
        IMAGES_DIR / "price_distribution.png",
        caption="Figure 1: Price distribution of men's jeans on Myntra",
    )
    add_image(
        doc,
        IMAGES_DIR / "price_segments.png",
        caption="Figure 2: Product count by price segment",
    )
    for item in [
        "Average selling price is ₹1,697 with a median of ₹1,484 — indicating a right-skewed distribution.",
        "50% of products are priced between ₹899 and ₹1,829 (affordable to mid-range segment).",
        "Premium products (up to ₹54,000) exist but represent a small portion of the catalog.",
        "Average customer rating is 3.98/5, with most products rated between 3.8 and 4.2.",
        "Roadster, HIGHLANDER, and United Colors of Benetton lead in total customer engagement.",
        "Premium products receive slightly higher ratings but lower overall engagement volume.",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    # Business Insights
    add_heading(doc, "5. Business Insights")

    insights = [
        (
            "5.1 Best Value for Money",
            "Top brands: MarvelQ, AngelFab, Metronaut, SZN, HARDSODA. These brands achieve a "
            "balanced combination of affordability, promotional offers, and customer satisfaction "
            "rather than relying on a single factor.",
        ),
        (
            "5.2 Customer Engagement",
            "HIGHLANDER leads in engagement efficiency (ratings per product). Roadster dominates "
            "total engagement due to its large catalog. Evaluating both metrics provides a "
            "comprehensive view of brand performance.",
        ),
        (
            "5.3 Promotional Discounts",
            "Some brands consistently rely on high average discounts. A small number of products "
            "receive 100% discounts, which may represent promotional anomalies. Discounts should "
            "be applied strategically to maintain profitability.",
        ),
        (
            "5.4 Price Segment Performance",
            "The mid-price segment (₹1,001–2,000) offers the largest product selection with "
            "strong engagement. Premium products achieve higher ratings but lower engagement, "
            "suggesting a smaller but satisfied customer base.",
        ),
        (
            "5.5 Brand Prioritization",
            "Top priority brands by Business Performance Score: HIGHLANDER, HARDSODA, Sztori, "
            "THE BEETEL HOUSE. These brands combine competitive pricing, positive ratings, "
            "effective promotions, and strong engagement.",
        ),
    ]
    insight_charts = {
        "5.1 Best Value for Money": IMAGES_DIR / "insights_023.png",
        "5.2 Customer Engagement": IMAGES_DIR / "insights_038.png",
        "5.4 Price Segment Performance": IMAGES_DIR / "insights_084.png",
        "5.5 Brand Prioritization": IMAGES_DIR / "insights_108.png",
    }
    for title, body in insights:
        add_heading(doc, title, level=2)
        doc.add_paragraph(body)
        chart = insight_charts.get(title)
        if chart:
            add_image(doc, chart)

    # Conclusions
    add_heading(doc, "6. Overall Conclusions")
    doc.add_paragraph(
        "The analysis reveals that Myntra's men's jeans catalog is concentrated in the mid-price "
        "segment, where customer engagement is highest. Brand success depends on a balanced "
        "combination of pricing, discounts, ratings, and engagement — not any single metric alone."
    )
    doc.add_paragraph(
        "Composite scoring models (Value Score and Business Performance Score) provide a more "
        "holistic evaluation than individual metrics, enabling Myntra to identify high-performing "
        "brands and optimize merchandising decisions."
    )

    # Recommendations
    add_heading(doc, "7. Recommendations")
    recommendations = [
        "Prioritize HIGHLANDER, HARDSODA, and Sztori for featured collections and marketing campaigns.",
        "Use value-for-money leaders (MarvelQ, AngelFab, Metronaut) as benchmarks for pricing strategy.",
        "Strengthen the mid-price portfolio (₹1,001–2,000) as the core revenue driver.",
        "Market premium products through exclusivity campaigns rather than discount-focused promotions.",
        "Monitor brands with consistently high discounts to ensure promotional profitability.",
        "Promote HIGHLANDER and Urbano Fashion for engagement efficiency despite smaller catalogs.",
        "Build an interactive dashboard for ongoing monitoring of brand and pricing metrics.",
    ]
    for rec in recommendations:
        doc.add_paragraph(rec, style="List Number")

    out = REPORTS_DIR / "Myntra_Data_Analysis_Report.docx"
    doc.save(out)
    print(f"Report saved: {out}")
    export_pdf(out)


def export_pdf(docx_path: Path) -> None:
    """Optional PDF export (Windows + Microsoft Word via docx2pdf)."""
    pdf_path = docx_path.with_suffix(".pdf")
    try:
        from docx2pdf import convert

        convert(str(docx_path), str(pdf_path))
        print(f"PDF saved: {pdf_path}")
    except ImportError:
        print("PDF export skipped (install docx2pdf: pip install docx2pdf)")
    except Exception as exc:
        print(f"PDF export skipped: {exc}")


def add_slide_title(prs, title, subtitle=None):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    left = PptInches(0.5)
    top = PptInches(2.5) if subtitle else PptInches(3)
    width = PptInches(9)

    box = slide.shapes.add_textbox(left, top, width, PptInches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptPt(32)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(0xE6, 0x2E, 0x5C)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        box2 = slide.shapes.add_textbox(left, PptInches(4), width, PptInches(1))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = PptPt(18)
        p2.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, bullets, image_path: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(14 if image_path else 16)
    if image_path and image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            PptInches(5.5),
            PptInches(1.8),
            width=PptInches(4.2),
        )
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    add_slide_title(
        prs,
        "Myntra Men's Jeans Data Analysis",
        "End-to-End Data Analytics Project | July 2026",
    )

    slides = [
        (
            "Problem Statement",
            [
                "Analyze Myntra's men's jeans catalog to uncover pricing and brand insights",
                "Support data-driven decisions for merchandising, marketing, and partnerships",
                "52,120 raw products → 31,527 cleaned records across 371 brands",
            ],
            None,
        ),
        (
            "Dataset Overview",
            [
                "7 features: brand, description, price, MRP, discount, ratings, # ratings",
                "Average price: ₹1,697 | Median: ₹1,484 | Range: ₹337 – ₹54,000",
                "Average rating: 3.98/5 | No missing values in raw data",
            ],
            None,
        ),
        (
            "Data Cleaning Process",
            [
                "Removed 17,047 exact duplicate rows",
                "Removed 3,546 records with inconsistent discount format (discount_percent > 1)",
                "Validated data types and discount format (decimal 0–1)",
                "Final clean dataset: 31,527 unique product records",
            ],
            None,
        ),
        (
            "EDA Highlights",
            [
                "Mid-range pricing dominates the catalog (50% between ₹899–₹1,829)",
                "Right-skewed price distribution — few luxury items pull mean up",
                "Roadster, HIGHLANDER lead in customer engagement",
            ],
            IMAGES_DIR / "price_distribution.png",
        ),
        (
            "BQ1: Best Value for Money",
            [
                "Composite Value Score = ratings + discount + price competitiveness",
                "Top brands: MarvelQ, AngelFab, Metronaut, SZN, HARDSODA",
                "Success comes from balanced pricing, not a single metric",
            ],
            IMAGES_DIR / "insights_023.png",
        ),
        (
            "BQ2: Customer Engagement",
            [
                "HIGHLANDER: highest engagement efficiency (ratings per product)",
                "Roadster: highest total engagement (large catalog)",
                "Evaluate both total and per-product metrics for full picture",
            ],
            IMAGES_DIR / "insights_038.png",
        ),
        (
            "BQ3: Promotional Discounts",
            [
                "Some brands rely heavily on high average discounts",
                "A few products show 100% discounts (promotional anomalies)",
                "Strategic discounting preserves profitability and brand value",
            ],
            None,
        ),
        (
            "BQ4: Price Segment Performance",
            [
                "Mid-price (₹1,001–2,000): largest catalog, strongest engagement",
                "Premium: higher ratings but fewer total ratings",
                "Core revenue driver = mid-range segment",
            ],
            IMAGES_DIR / "price_segments.png",
        ),
        (
            "BQ5: Brand Prioritization",
            [
                "Business Performance Score integrates all key metrics",
                "Top brands: HIGHLANDER, HARDSODA, Sztori, THE BEETEL HOUSE",
                "Prioritize for featured collections and inventory expansion",
            ],
            IMAGES_DIR / "insights_108.png",
        ),
        (
            "Final Recommendations",
            [
                "Feature HIGHLANDER, HARDSODA, Sztori in campaigns",
                "Strengthen mid-price portfolio as core revenue segment",
                "Market premium via exclusivity, not discounts",
                "Use Streamlit dashboard for ongoing brand monitoring",
            ],
            None,
        ),
    ]

    for title, bullets, image in slides:
        add_content_slide(prs, title, bullets, image)

    # Thank you slide
    add_slide_title(prs, "Thank You", "Questions & Discussion")

    out = REPORTS_DIR / "Myntra_Data_Analysis_Presentation.pptx"
    prs.save(out)
    print(f"Presentation saved: {out}")


if __name__ == "__main__":
    build_report()
    build_presentation()
