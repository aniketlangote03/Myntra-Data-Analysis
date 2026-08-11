import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
REPORTS_DIR = ROOT / "reports"
IMAGES_DIR = ROOT / "images"
REPORTS_DIR.mkdir(exist_ok=True)

# Color Palette Definitions
NAVY_DARK = RGBColor(0x0F, 0x17, 0x2A)     # #0F172A - Header & Primary Text
MYNTRA_CORAL = RGBColor(0xE6, 0x2E, 0x5C)  # #E62E5C - Accent / Myntra Brand
PBI_YELLOW = RGBColor(0xF5, 0x9E, 0x0B)    # #F59E0B - Power BI Accent
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)      # #F8FAFC - Card Background
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)    # #64748B - Secondary Subtext
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def create_deck():
    prs = Presentation()
    # Set 16:9 Widescreen Dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Helper: Add Slide Background
    def set_bg(slide, color=BG_LIGHT):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    # Helper: Add Standard Header
    def add_header(slide, title_text, category_text="POWER BI DATA ANALYST PRESENTATION"):
        # Header background bar
        hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = NAVY_DARK
        hbar.line.fill.background()

        # Category tag
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(10), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = MYNTRA_CORAL

        # Slide Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

    # Helper: Add Card
    def add_card(slide, left, top, width, height, title, items, border_color=MYNTRA_CORAL):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = f"• {item}"
            p_item.font.size = Pt(13)
            p_item.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
            p_item.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, NAVY_DARK)

    # Title Card Accent Box
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(0.15), Inches(3.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = MYNTRA_CORAL
    accent_bar.line.fill.background()

    tb = s1.shapes.add_textbox(Inches(1.4), Inches(2.0), Inches(10.5), Inches(3.5))
    tf = tb.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Myntra Men's Jeans Analytics"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "Executive Power BI & End-to-End Data Analytics Presentation"
    p2.font.size = Pt(20)
    p2.font.color.rgb = PBI_YELLOW
    p2.space_after = Pt(30)

    p3 = tf.add_paragraph()
    p3.text = "Prepared by: Aniket Langote  |  Role: Data Analyst\nTools: Power BI Desktop • DAX • Python (Pandas/Plotly) • Streamlit\nDataset: 31,527 Cleaned E-Commerce Listings across 371 Brands"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary & Project Context
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "Executive Summary & Business Context")

    # 3 Summary Column Cards
    col_w = Inches(3.7)
    gap = Inches(0.3)
    t_top = Inches(1.6)
    h_card = Inches(5.2)

    add_card(s2, Inches(0.8), t_top, col_w, h_card, "1. Business Objective", [
        "Optimize e-commerce catalog pricing and discount depth on Myntra.",
        "Identify high-performing brand partners for platform growth.",
        "Evaluate customer rating trends vs review engagement volume.",
        "Deliver scalable BI deliverables for leadership decision-making."
    ], MYNTRA_CORAL)

    add_card(s2, Inches(0.8) + col_w + gap, t_top, col_w, h_card, "2. Data Scale & Scope", [
        "Raw scraped dataset: 52,120 product listings.",
        "Cleaned & validated dataset: 31,527 records across 371 brands.",
        "Features: Brand Name, Description, Price (₹), MRP (₹), Discount %, Ratings, Review Count.",
        "Total customer review volume analyzed: 3.31 Million ratings."
    ], PBI_YELLOW)

    add_card(s2, Inches(0.8) + (col_w + gap)*2, t_top, col_w, h_card, "3. Multi-Channel Analytics", [
        "Power BI Desktop Dashboard (.pbix): Interactive DAX modeling & visual reporting.",
        "Streamlit Web Application: Real-time dynamic filtering & brand scoring.",
        "Plotly HTML Dashboard: Client-side standalone web dashboard.",
        "Automated Reports: Word (.docx), PDF export, and PowerPoint (.pptx)."
    ], NAVY_DARK)

    # -------------------------------------------------------------
    # SLIDE 3: Data Hygiene & ETL Pipeline
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "Data Cleaning & Quality Assurance Pipeline")

    # Flowchart-like Horizontal Steps
    step_w = Inches(2.7)
    step_gap = Inches(0.3)
    s_top = Inches(1.8)
    s_h = Inches(4.8)

    steps = [
        ("Step 1: Raw Ingestion", "52,120 Raw Rows", [
            "Loaded scraped CSV catalog.",
            "Evaluated 7 core schema columns.",
            "Assessed missing values (0 missing found)."
        ], RGBColor(0x3B, 0x82, 0xF6)),
        
        ("Step 2: Deduplication", "-17,047 Duplicate Rows", [
            "Identified exact row duplicates due to scraper overlap.",
            "Reduced volume from 52.1k → 35,073 rows.",
            "Preserved single unique records."
        ], RGBColor(0x10, 0xB9, 0x81)),
        
        ("Step 3: Quality Cleansing", "-3,546 Format Anomalies", [
            "Stripped invalid discounts (>100% / scraper scale errors).",
            "Cast numerical types & validated pricing bounds.",
            "Final clean dataset: 31,527 items."
        ], MYNTRA_CORAL),
        
        ("Step 4: DAX Data Model", "31,527 Clean Items", [
            "Imported clean dataset to Power BI.",
            "Built DAX calculated metrics.",
            "Engineered Value & Business Performance scores."
        ], NAVY_DARK)
    ]

    for idx, (stitle, sbadge, sbullets, scolor) in enumerate(steps):
        s_left = Inches(0.8) + idx * (step_w + step_gap)
        add_card(s3, s_left, s_top, step_w, s_h, stitle, sbullets, scolor)
        
        # Badge pill at top of card
        bpill = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_left + Inches(0.2), s_top + Inches(0.6), step_w - Inches(0.4), Inches(0.35))
        bpill.fill.solid()
        bpill.fill.fore_color.rgb = scolor
        bpill.line.fill.background()
        tf_b = bpill.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = sbadge
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE
        p_b.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 4: Power BI Architecture & DAX Modeling
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "Power BI Data Model & Calculated DAX Measures")

    # Left Column: DAX Measures Table / Box
    add_card(s4, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2), "Key DAX Calculated Measures", [
        "Total Products = COUNTROWS('myntra_cleaned') → 31,527",
        "Avg Selling Price = AVERAGE('myntra_cleaned'[price]) → ₹1,697",
        "Avg Discount % = AVERAGE('myntra_cleaned'[discount_percent]) → 50.3%",
        "Total Reviews = SUM('myntra_cleaned'[number_of_ratings]) → 3.31 Million",
        "Weighted Rating = SUMX(Table, Rating * Reviews) / Total Reviews",
        "Value Score Index = 0.40*Rating_Norm + 0.30*Disc_Norm + 0.30*Price_Affordability",
        "Business Performance Score = 0.35*Rating + 0.30*Engagement + 0.20*Value + 0.15*Discount"
    ], MYNTRA_CORAL)

    # Right Column: Dashboard Design Features
    add_card(s4, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2), "Power BI Technical Highlights", [
        "Executive Color Theme: Deep Navy (#0B2B6F), Coral Accent, & White Card Visuals.",
        "Dynamic Report Slicers: Interactive filters by Brand Name, Price Segment, and Rating Tier.",
        "Custom Icon KPI Cards: Custom PNG asset integration for star, discount, price, and customer icons.",
        "4K Widescreen Canvas: 4000 x 2250 responsive executive reporting canvas.",
        "Interactive Cross-Filtering: Real-time slicer interaction across all visual panels."
    ], PBI_YELLOW)

    # -------------------------------------------------------------
    # SLIDE 5: Power BI Executive Dashboard Overview (Embedded Screenshot)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "Power BI Executive Dashboard Live Tour")

    # Embed Power BI Dashboard Image
    pbi_img = IMAGES_DIR / "powerbi_dashboard_v4.png"
    if pbi_img.exists():
        s5.shapes.add_picture(str(pbi_img), Inches(0.8), Inches(1.5), width=Inches(11.73))

    # Caption Box below image
    c_box = s5.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11.73), Inches(0.5))
    tf_c = c_box.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = "Power BI Dashboard File: PowerBi/Myntra Sales Analytics.pbix  |  Featuring Header Slicers, 6 KPI Cards & 10 Interactive Visual Panels"
    p_c.font.size = Pt(11)
    p_c.font.color.rgb = TEXT_MUTED
    p_c.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 6: Key Insight 1 - Pricing Strategy & Catalog Mix
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "Insight 1: Price Distribution & Catalog Mix")

    add_card(s6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Catalog Pricing Highlights", [
        "Average Selling Price: ₹1,697 | Median Price: ₹1,484.",
        "Right-Skewed Distribution: Luxury items (up to ₹54,000) elevate the average price above median.",
        "Core Revenue Segment: Mid-Range (₹1,000–₹2,000) represents >50% of total product volume.",
        "Promotional Discount Depth: Average discount is ~50.3%, with mass-market brands reaching 65%+ discount reliance.",
        "Price Bands: Budget (<₹1k), Mid-Range (₹1k-2k), Premium (₹2k-3.5k), Luxury (>₹3.5k)."
    ], MYNTRA_CORAL)

    # Embed Price Distribution Chart
    chart_price = IMAGES_DIR / "price_distribution.png"
    if chart_price.exists():
        s6.shapes.add_picture(str(chart_price), Inches(6.7), Inches(1.8), width=Inches(5.8))

    # -------------------------------------------------------------
    # SLIDE 7: Key Insight 2 - Brand Volume vs Engagement
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "Insight 2: Brand Catalog Share vs Customer Engagement")

    add_card(s7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Brand Performance Dynamics", [
        "Volume Leaders: United Colors of Benetton (2,992 items), Flying Machine (2,575 items), Roadster (1,794 items).",
        "Total Engagement Champion: Roadster dominates total customer reviews with 532,000+ ratings.",
        "Engagement Efficiency: HIGHLANDER leads in ratings per product, indicating high customer conversion.",
        "Rating Benchmark: Catalog average rating is 3.98/5; premium brands achieve slightly higher ratings but lower aggregate volume."
    ], PBI_YELLOW)

    # Embed Brand Engagement Chart
    chart_eng = IMAGES_DIR / "top_brands_engagement.png"
    if chart_eng.exists():
        s7.shapes.add_picture(str(chart_eng), Inches(6.7), Inches(1.8), width=Inches(5.8))

    # -------------------------------------------------------------
    # SLIDE 8: Key Insight 3 - Value Score & Business Performance
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "Insight 3: Composite Value & Business Performance Scores")

    add_card(s8, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), "Composite Ranking Models", [
        "Value Score Model: Combines Customer Rating (40%), Discount Depth (30%), and Price Affordability (30%).",
        "Top Value Champions: HARDSODA, MarvelQ, AngelFab, Metronaut deliver optimal quality-to-price ratio.",
        "Business Performance Score (BPS): Evaluates Rating (35%), Review Efficiency (30%), Value Score (20%), and Discounting (15%).",
        "Top Priority Brands: HIGHLANDER (164.59), Roadster (131.23), Flying Machine (107.52), Benetton (114.67)."
    ], NAVY_DARK)

    # Embed Business Performance Chart
    chart_bps = IMAGES_DIR / "insights_108.png"
    if chart_bps.exists():
        s8.shapes.add_picture(str(chart_bps), Inches(6.7), Inches(1.8), width=Inches(5.8))

    # -------------------------------------------------------------
    # SLIDE 9: Strategic Data Analyst Recommendations
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_header(s9, "Strategic Recommendations for Myntra Merchandising")

    rec_w = Inches(5.6)
    rec_h = Inches(2.4)

    # 4 Quadrant Cards
    add_card(s9, Inches(0.8), Inches(1.6), rec_w, rec_h, "1. Merchandising & Campaign Focus", [
        "Prioritize HIGHLANDER, HARDSODA, and Roadster for homepage banners and featured sales.",
        "Expand catalog listings for high-engagement, high-value score brands."
    ], MYNTRA_CORAL)

    add_card(s9, Inches(6.9), Inches(1.6), rec_w, rec_h, "2. Pricing Tier Optimization", [
        "Strengthen the core ₹1,000–₹2,000 mid-price segment as primary revenue driver.",
        "Establish value-for-money benchmarks (HARDSODA/MarvelQ) for new brand onboarding."
    ], PBI_YELLOW)

    add_card(s9, Inches(0.8), Inches(4.3), rec_w, rec_h, "3. Discount & Margin Protection", [
        "Monitor brands relying heavily on >60% discounts to protect gross margin.",
        "Promote premium lines through brand exclusivity rather than steep discount erosion."
    ], NAVY_DARK)

    add_card(s9, Inches(6.9), Inches(4.3), rec_w, rec_h, "4. Interactive BI Monitoring", [
        "Deploy Power BI Dashboard to category managers for real-time brand filtering.",
        "Track customer rating trends dynamically to detect catalog quality shifts."
    ], RGBColor(0x10, 0xB9, 0x81))

    # -------------------------------------------------------------
    # SLIDE 10: Conclusion & Q&A Slide (Dark Theme)
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, NAVY_DARK)

    tb10 = s10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf10 = tb10.text_frame

    p10_1 = tf10.paragraphs[0]
    p10_1.text = "Thank You!"
    p10_1.font.size = Pt(44)
    p10_1.font.bold = True
    p10_1.font.color.rgb = MYNTRA_CORAL
    p10_1.alignment = PP_ALIGN.CENTER
    p10_1.space_after = Pt(15)

    p10_2 = tf10.add_paragraph()
    p10_2.text = "Myntra Men's Jeans Data Analytics & Power BI Presentation"
    p10_2.font.size = Pt(20)
    p10_2.font.bold = True
    p10_2.font.color.rgb = WHITE
    p10_2.alignment = PP_ALIGN.CENTER
    p10_2.space_after = Pt(25)

    p10_3 = tf10.add_paragraph()
    p10_3.text = "Project Deliverables in Repository:\n• Power BI File: PowerBi/Myntra Sales Analytics.pbix\n• Streamlit Dashboard: dashboard/app.py\n• Plotly HTML Dashboard: dashboard/index.html\n• Python Notebooks: notebooks/01_Data_Loading.ipynb to 05_Advanced_Business_Insights.ipynb\n• Executive Word/PDF Reports: reports/Myntra_Data_Analysis_Report.docx & .pdf"
    p10_3.font.size = Pt(13)
    p10_3.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    p10_3.alignment = PP_ALIGN.CENTER

    # Save output deck
    out_pptx = REPORTS_DIR / "Myntra_PowerBI_Data_Analyst_Presentation.pptx"
    prs.save(out_pptx)
    print(f"Data Analyst PowerPoint Presentation saved to {out_pptx}")

if __name__ == '__main__':
    create_deck()
